import os
import sys
import pickle
import time
from dotenv import load_dotenv

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_ollama import OllamaEmbeddings
from langchain_community.vectorstores import Chroma
from langchain_community.retrievers import BM25Retriever

def sanitise_text(text: str) -> str:
    """Sanitise text content to protect against prompt injection, safety override tricks, and formatting errors."""
    if not text:
        return ""
    import re
    patterns_to_neutralize = [
        r"ignore\s+(?:all\s+)?previous\s+instructions",
        r"ignore\s+(?:any\s+)?instructions\s+above",
        r"ignore\s+the\s+(?:context|rules|guidelines)",
        r"system\s+override",
        r"you\s+are\s+now\s+a",
        r"instead\s+of\s+answering",
        r"forget\s+your\s+previous",
        r"bypass\s+safety",
        r"delete\s+all",
        r"reset\s+database",
        r"execute\s+command",
    ]
    sanitised = text
    for pattern in patterns_to_neutralize:
        sanitised = re.sub(pattern, "[CLEANED SECURE DATA]", sanitised, flags=re.IGNORECASE)
    sanitised = sanitised.replace("{", "[").replace("}", "]")
    sanitised = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", sanitised)
    return sanitised

def ingest_document(file_path, session_id="default"):
    ext = os.path.splitext(file_path)[1].lower()
    print(f"Loading document: {file_path} (type: {ext}) in session {session_id}...")
    
    from langchain_core.documents import Document
    docs = []
    
    easyocr_reader = None
    
    if ext == ".pdf":
        import pypdf
        reader = pypdf.PdfReader(file_path)
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip() and len(text.strip()) > 30:
                print(f"Extracted direct text from page {i + 1}")
                docs.append(Document(
                    page_content=sanitise_text(text), 
                    metadata={"source": file_path, "page": i + 1}
                ))
            else:
                print(f"Page {i + 1} appears to be scanned. Running OCR...")
                try:
                    import fitz  # PyMuPDF
                    import tempfile
                    
                    doc = fitz.open(file_path)
                    fitz_page = doc.load_page(i)
                    pix = fitz_page.get_pixmap()
                    
                    with tempfile.TemporaryDirectory() as temp_dir:
                        image_path = os.path.join(temp_dir, f"temp_page_{i}.png")
                        pix.save(image_path)
                        
                        if easyocr_reader is None:
                            import easyocr
                            print("Initializing EasyOCR reader...")
                            easyocr_reader = easyocr.Reader(['en'], gpu=False)
                            
                        result = easyocr_reader.readtext(image_path, detail=0)
                        ocr_text = "\n".join(result)
                        if ocr_text.strip():
                            docs.append(Document(
                                page_content=sanitise_text(ocr_text), 
                                metadata={"source": file_path, "page": i + 1, "ocr": True}
                            ))
                except Exception as e:
                    print(f"OCR failed for page {i + 1}: {e}")
                    
    elif ext in (".png", ".jpg", ".jpeg"):
        try:
            import easyocr
            print("Initializing EasyOCR reader...")
            easyocr_reader = easyocr.Reader(['en'], gpu=False)
            print(f"Running OCR on image {file_path}...")
            result = easyocr_reader.readtext(file_path, detail=0)
            ocr_text = "\n".join(result)
            if ocr_text.strip():
                docs.append(Document(
                    page_content=sanitise_text(ocr_text), 
                    metadata={"source": file_path, "ocr": True}
                ))
        except Exception as e:
            raise ValueError(f"Failed to run OCR on image: {e}")
            
    elif ext == ".docx":
        from langchain_community.document_loaders import Docx2txtLoader
        loader = Docx2txtLoader(file_path)
        raw_docs = loader.load()
        for d in raw_docs:
            docs.append(Document(page_content=sanitise_text(d.page_content), metadata=d.metadata))
            
    elif ext == ".doc":
        try:
            from langchain_community.document_loaders import Docx2txtLoader
            loader = Docx2txtLoader(file_path)
            raw_docs = loader.load()
            for d in raw_docs:
                docs.append(Document(page_content=sanitise_text(d.page_content), metadata=d.metadata))
        except Exception:
            raise ValueError("Word .doc format is not supported directly. Please save your file as .docx and try again.")
            
    elif ext in (".md", ".txt"):
        from langchain_community.document_loaders import TextLoader
        loader = TextLoader(file_path, encoding='utf-8')
        raw_docs = loader.load()
        for d in raw_docs:
            docs.append(Document(page_content=sanitise_text(d.page_content), metadata=d.metadata))
            
    elif ext == ".csv":
        import csv
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.DictReader(f)
            for row_idx, row in enumerate(reader):
                row_text = ", ".join([f"{k}: {v}" for k, v in row.items() if v])
                if row_text.strip():
                    docs.append(Document(
                        page_content=sanitise_text(row_text), 
                        metadata={"source": file_path, "row": row_idx}
                    ))
                    
    elif ext == ".xlsx":
        import pandas as pd
        xl = pd.ExcelFile(file_path)
        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name)
            df = df.fillna("")
            for row_idx, row in df.iterrows():
                row_text = ", ".join([f"{col}: {val}" for col, val in row.items() if str(val).strip()])
                if row_text.strip():
                    docs.append(Document(
                        page_content=sanitise_text(row_text), 
                        metadata={"source": file_path, "sheet": sheet_name, "row": row_idx}
                    ))

    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                slide_text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text_parts.append(shape.text.strip())
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    slide_text_parts.append(cell.text.strip())
                slide_text = "\n".join(slide_text_parts)
                if slide_text.strip():
                    docs.append(Document(
                        page_content=sanitise_text(slide_text),
                        metadata={"source": file_path, "slide": i + 1}
                    ))
        except Exception as e:
            raise ValueError(f"Failed to parse PowerPoint presentation: {e}")
    else:
        raise ValueError(f"Unsupported file type: {ext}")
        
    if not docs:
        raise ValueError(f"No text content could be extracted from {file_path}")
        
    # Assign doc_id to every loaded document
    doc_id = os.path.basename(file_path)
    for d in docs:
        d.metadata["doc_id"] = doc_id
        d.metadata["session_id"] = session_id
        
    print(f"Loaded {len(docs)} document section(s) for doc_id: {doc_id} in session {session_id}.")
    
    print("Splitting text into chunks...")
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=80,
        separators=["\n\n", "\n", " ", ""]
    )
    splits = text_splitter.split_documents(docs)
    
    # Filter out empty splits
    splits = [s for s in splits if s.page_content.strip()]
    for s in splits:
        s.metadata["doc_id"] = doc_id
        s.metadata["session_id"] = session_id
    print(f"Created {len(splits)} chunks.")

    print("Generating Vector Embeddings and saving to Chroma in batches...")
    chroma_db_path = os.getenv("CHROMA_DB_PATH", "./chroma_db")
    ollama_embed_model = os.getenv("OLLAMA_EMBED_MODEL", "nomic-embed-text")

    embeddings = OllamaEmbeddings(model=ollama_embed_model)

    # Use existing Chroma database, clear old version of this document if it exists, and append new chunks.
    vectorstore = Chroma(persist_directory=chroma_db_path, embedding_function=embeddings)
    try:
        print(f"Checking for existing index entries for {doc_id} in session {session_id} to prevent duplication...")
        vectorstore.delete(where={"$and": [{"doc_id": doc_id}, {"session_id": session_id}]})
        print(f"Cleared existing chunks for {doc_id} in session {session_id}.")
    except Exception as e:
        print(f"No previous index entry to clean or collection is empty: {e}")
    
    batch_size = 50
    for i in range(0, len(splits), batch_size):
        batch = splits[i:i + batch_size]
        print(f"Processing batch {i//batch_size + 1} (chunks {i} to {min(i+batch_size, len(splits))})...")
        vectorstore.add_documents(documents=batch)
        
    print(f"Vectors saved to {chroma_db_path}")
    print("Ingestion complete!")

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Ingest document into RAG database")
    parser.add_argument("file_path", help="Path to the file to ingest")
    parser.add_argument("--session-id", default="default", help="Private session ID for chat-level isolation")
    args = parser.parse_args()
    
    if not os.path.exists(args.file_path):
        print(f"Error: File {args.file_path} does not exist.")
        sys.exit(1)
        
    ingest_document(args.file_path, session_id=args.session_id)
